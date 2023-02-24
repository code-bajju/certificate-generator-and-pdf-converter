import openpyxl
from pptx import Presentation
import comtypes.client
import os

# Define function to convert PPTX to PDF
def pptx_to_pdf(input_file, output_file):
    # Get the absolute paths of the input and output files
    input_file = os.path.abspath(input_file)
    output_file = os.path.abspath(output_file)

    # Create a PowerPoint application object
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")

    # Set visibility to False so the application doesn't open a new window
    powerpoint.Visible = True

    # Open the input file
    deck = powerpoint.Presentations.Open(input_file)

    # Save the deck as PDF
    deck.SaveAs(output_file, 32)

    # Close the PowerPoint application
    powerpoint.Quit()


def update_certificates():
    # Load the Excel file containing the new names
    wb = openpyxl.load_workbook('names.xlsx')
    ws = wb.active

    # Load the PowerPoint template
    template = 'Certificate.ppt'

    # Loop through each row in the Excel file
    for i in range(1, 36):
        # Get the new name from the Excel file
        new_name = ws.cell(row=i, column=1).value
        # Load the PowerPoint presentation
        prs = Presentation(template)
        # Loop through each slide in the presentation
        for slide in prs.slides:
            # Loop through each shape in the slide
            for shape in slide.shapes:
                # Check if the shape is a text box
                if shape.has_text_frame:
                    # Get the text frame
                    text_frame = shape.text_frame
                    # Loop through each paragraph in the text frame
                    for paragraph in text_frame.paragraphs:
                        # Loop through each run in the paragraph
                        for run in paragraph.runs:
                            # Get the text from the run
                            old_text = run.text
                            # Check if the old text matches "Subham Kumar"
                            if old_text == "Subham Kumar":
                                # Update the text with the new name
                                run.text = new_name
        # Save the updated PowerPoint presentation with a new file name
        prs.save(f'{new_name}_certificate.ppt')
        # Convert the updated PowerPoint presentation to PDF format
        pptx_to_pdf(f'{new_name}_certificate.ppt', f'{new_name}_certificate.pdf')

    # Close the Excel file
    wb.close()


# Create a menu to select options
while True:
    print("1. Update certificates")
    print("2. Exit")
    choice = input("Enter your choice (1/2): ")

    if choice == "1":
        update_certificates()
        print("Certificates updated successfully!")
    elif choice == "2":
        break
    else:
        print("Invalid choice. Please try again.")
